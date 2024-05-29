import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import os
import japanize_matplotlib
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from PIL import Image
import streamlit as st
from io import BytesIO

Image.MAX_IMAGE_PIXELS = None

prs = Presentation()

def load_excel(path):
    df = pd.ExcelFile(path)

    return df


def clean_sheet(df, sheet_name):
    df_1 = pd.read_excel(df, sheet_name, header=None)
    df_2 = df_1.dropna(axis=1, how='all', subset=[0])
    df_3 = df_2.iloc[:, :5]
    df_4 = df_3.dropna(how='all')
    df_5 = df_4.iloc[1:, :]
    df_5.reset_index(inplace=True, drop=True)
    df_6 = df_5.fillna('')
    df_6[1] = df_6[1].str.strip()

    return df_6


def create_picture(df, target_rows_lb, target_rows_db, name):
    colors = []

    for row in range(0, df.shape[0]):
        if row in target_rows_lb:
            colors.append(["#c1c9db","#c1c9db","#c1c9db","#c1c9db","#c1c9db"])
        elif row in target_rows_db:
            colors.append(["#8191b5","#8191b5","#8191b5","#8191b5","#8191b5"])
        else:
            colors.append(["w","w","w","w","w"])

    fig, ax = plt.subplots()

    ax.axis('tight')
    ax.axis('off')

    plt.rcParams['text.color'] = 'black'

    table = ax.table(cellText=df.values, cellColours=colors,  loc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1.5,2)

    for c in table.get_children():
        c.set_edgecolor('none')

    fullname = 'mabpfolder/' + name

    plt.savefig(fullname, bbox_inches='tight',pad_inches=0.6, dpi=1000)

    return fullname


def create_slide_1(df_6):
    global prs

    # picture 1
    name_list1 = df_6.index[df_6[1] == '資産合計'].tolist()
    df_p1 = df_6.iloc[:name_list1[0]+1]
    target_rows_lb1 = []
    target_rows_db1 = [1]
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '流動資産'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '固定資産'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '資産合計'].tolist()
    pic_p1 = create_picture(df_p1, target_rows_lb1, target_rows_db1, 's1p1.png')

    # picture 2
    name_list2 = df_6.index[df_6[1] == '負債合計'].tolist() 
    df_p2 = df_6.iloc[name_list1[0]+1:name_list2[0]+1]
    target_rows_lb2 = []
    target_rows_db2 = [1]
    target_rows_lb2 = target_rows_lb2 + df_p2.index[df_p2[1] == '流動負債'].tolist()
    target_rows_lb2 = target_rows_lb2 + df_p2.index[df_p2[1] == '固定負債'].tolist()
    target_rows_db2 = target_rows_db2 + df_p2.index[df_p2[1] == '負債合計'].tolist()
    pic_p2 = create_picture(df_p2, target_rows_lb2, target_rows_db2, 's1p2.png')

    # picture 3
    name_list3 = df_6.index[df_6[1] == '負債・純資産合計'].tolist() 
    df_p3 = df_6.iloc[name_list2[0]+1:name_list3[0]+1]
    target_rows_lb3 = []
    target_rows_db3 = []
    target_rows_lb3 = target_rows_lb3 + df_p3.index[df_p3[1] == '株主資本'].tolist()
    target_rows_db3 = target_rows_lb3 + df_p3.index[df_p3[1] == '純資産合計'].tolist()
    target_rows_db3 = target_rows_db3 + df_p3.index[df_p3[1] == '負債・純資産合計'].tolist()
    pic_p3 = create_picture(df_p3, target_rows_lb3, target_rows_db3, 's1p3.png')

    # scale picture 1
    img1 = Image.open(pic_p1)
    original_width1, original_height1 = img1.size
    new_height1 = 6
    scaling_factor1 = new_height1/(original_height1/1000)
    new_width1 = scaling_factor1*original_width1/1000

    # scale picture 2
    img2 = Image.open(pic_p2)
    original_width2, original_height2 = img2.size
    new_height2 = 3
    scaling_factor2 = new_height2/(original_height2/1000)
    new_width2 = scaling_factor2*original_width2/1000

    # scale picture 3
    img3 = Image.open(pic_p3)
    original_width3, original_height3 = img3.size
    new_height3 = 2
    scaling_factor3 = new_height3/(original_height3/1000)
    new_width3 = scaling_factor3*original_width3/1000

    # presentation initialize
    slide_layout1 = prs.slide_layouts[5]  # Using a title-only layout
    slide1 = prs.slides.add_slide(slide_layout1)
    title = slide1.shapes.title
    title.text = "貸借対照表（3期比較）"
    if title.text_frame:
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
    pic1 = slide1.shapes.add_picture(pic_p1, Inches(0.5), Inches(1.3), width=Inches(new_width1), height=Inches(new_height1))
    pic2 = slide1.shapes.add_picture(pic_p2, Inches(4.5), Inches(1.3), width=Inches(new_width2), height=Inches(new_height2))
    pic3 = slide1.shapes.add_picture(pic_p3, Inches(4.5), Inches(4.3), width=Inches(new_width3), height=Inches(new_height3))


def create_slide_1(df_6):
    global prs

    # picture 1
    name_list1 = df_6.index[df_6[1] == '資産合計'].tolist()
    df_p1 = df_6.iloc[:name_list1[0]+1]
    target_rows_lb1 = []
    target_rows_db1 = [1]
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '流動資産'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '固定資産'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '資産合計'].tolist()
    pic_p1 = create_picture(df_p1, target_rows_lb1, target_rows_db1, 's1p1.png')

    # picture 2
    name_list2 = df_6.index[df_6[1] == '負債合計'].tolist() 
    df_p2 = df_6.iloc[name_list1[0]+1:name_list2[0]+1]
    target_rows_lb2 = []
    target_rows_db2 = [1]
    target_rows_lb2 = target_rows_lb2 + df_p2.index[df_p2[1] == '流動負債'].tolist()
    target_rows_lb2 = target_rows_lb2 + df_p2.index[df_p2[1] == '固定負債'].tolist()
    target_rows_db2 = target_rows_db2 + df_p2.index[df_p2[1] == '負債合計'].tolist()
    pic_p2 = create_picture(df_p2, target_rows_lb2, target_rows_db2, 's1p2.png')

    # picture 3
    name_list3 = df_6.index[df_6[1] == '負債・純資産合計'].tolist() 
    df_p3 = df_6.iloc[name_list2[0]+1:name_list3[0]+1]
    target_rows_lb3 = []
    target_rows_db3 = []
    target_rows_lb3 = target_rows_lb3 + df_p3.index[df_p3[1] == '株主資本'].tolist()
    target_rows_db3 = target_rows_lb3 + df_p3.index[df_p3[1] == '純資産合計'].tolist()
    target_rows_db3 = target_rows_db3 + df_p3.index[df_p3[1] == '負債・純資産合計'].tolist()
    pic_p3 = create_picture(df_p3, target_rows_lb3, target_rows_db3, 's1p3.png')

    # scale picture 1
    img1 = Image.open(pic_p1)
    original_width1, original_height1 = img1.size
    new_height1 = 6
    scaling_factor1 = new_height1/(original_height1/1000)
    new_width1 = scaling_factor1*original_width1/1000

    # scale picture 2
    img2 = Image.open(pic_p2)
    original_width2, original_height2 = img2.size
    new_height2 = 3
    scaling_factor2 = new_height2/(original_height2/1000)
    new_width2 = scaling_factor2*original_width2/1000

    # scale picture 3
    img3 = Image.open(pic_p3)
    original_width3, original_height3 = img3.size
    new_height3 = 2
    scaling_factor3 = new_height3/(original_height3/1000)
    new_width3 = scaling_factor3*original_width3/1000

    # presentation initialize
    slide_layout1 = prs.slide_layouts[5]  # Using a title-only layout
    slide1 = prs.slides.add_slide(slide_layout1)
    title = slide1.shapes.title
    title.text = "貸借対照表（3期比較）"
    if title.text_frame:
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
    pic1 = slide1.shapes.add_picture(pic_p1, Inches(0.5), Inches(1.3), width=Inches(new_width1), height=Inches(new_height1))
    pic2 = slide1.shapes.add_picture(pic_p2, Inches(4.5), Inches(1.3), width=Inches(new_width2), height=Inches(new_height2))
    pic3 = slide1.shapes.add_picture(pic_p3, Inches(4.5), Inches(4.3), width=Inches(new_width3), height=Inches(new_height3))


def create_slide_2(df_6):
    global prs

    # picture 1
    name_list1 = df_6.index[df_6[1] == '当期純利益'].tolist()
    df_p1 = df_6.iloc[:name_list1[0]+1]
    target_rows_lb1 = []
    target_rows_db1 = [1]
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '売上高'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '売上原価'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '販売費及び一般管理費'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '営業外収益'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '営業外費用'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '特別利益'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '特別損失'].tolist()
    target_rows_lb1 = target_rows_lb1 + df_p1.index[df_p1[1] == '法人税等'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '売上総利益'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '営業利益'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '経常利益'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '税引前当期純利益'].tolist()
    target_rows_db1 = target_rows_db1 + df_p1.index[df_p1[1] == '当期純利益'].tolist()
    pic_p1 = create_picture(df_p1, target_rows_lb1, target_rows_db1, 's2p1.png')

    # picture 2
    name_list2_1 = df_6.index[df_6[1] == '製造原価（3ヵ年推移）'].tolist() 
    name_list2_2 = df_6.index[df_6[1] == '合計'].tolist()
    df_p2 = df_6.iloc[name_list2_1[0]:name_list2_2[0]+1]
    target_rows_lb2 = [2]
    target_rows_db2 = [1, name_list2_2]
    pic_p2 = create_picture(df_p2, target_rows_lb2, target_rows_db2, 's1p2.png')

    # scale picture 1
    img1 = Image.open(pic_p1)
    original_width1, original_height1 = img1.size
    new_height1 = 6
    scaling_factor1 = new_height1/(original_height1/1000)
    new_width1 = scaling_factor1*original_width1/1000

    # scale picture 2
    img2 = Image.open(pic_p2)
    original_width2, original_height2 = img2.size
    new_height2 = 5
    scaling_factor2 = new_height2/(original_height2/1000)
    new_width2 = scaling_factor2*original_width2/1000

    # presentation initialize
    slide_layout2 = prs.slide_layouts[5]  # Using a title-only layout
    slide2 = prs.slides.add_slide(slide_layout2)
    title = slide2.shapes.title
    title.text = "損益計算書（3期比較）"
    if title.text_frame:
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
    pic1 = slide2.shapes.add_picture(pic_p1, Inches(0.5), Inches(1.3), width=Inches(new_width1), height=Inches(new_height1))
    pic2 = slide2.shapes.add_picture(pic_p2, Inches(4.5), Inches(1.3), width=Inches(new_width2), height=Inches(new_height2))


def create_slide_3(df_6):
    global prs

    # picture 1
    name_list1_1 = df_6.index[df_6[1] == '製造原価（3ヵ年推移）'].tolist()
    name_list1_2 = df_6.index[df_6[1] == '合計'].tolist()
    df_p1 = df_6.iloc[name_list1_1[0]+1:name_list1_2[0]+1]
    target_rows_lb1 = [name_list1_2]
    target_rows_db1 = [1, 2]
    pic_p1 = create_picture(df_p1, target_rows_lb1, target_rows_db1, 's3p1.png')

    # scale picture 1
    img1 = Image.open(pic_p1)
    original_width1, original_height1 = img1.size
    new_height1 = 6
    scaling_factor1 = new_height1/(original_height1/1000)
    new_width1 = scaling_factor1*original_width1/1000

    # presentation initialize
    slide_layout3 = prs.slide_layouts[5]  # Using a title-only layout
    slide3 = prs.slides.add_slide(slide_layout3)
    title = slide3.shapes.title
    title.text = "損益計算書（3期比較）"
    if title.text_frame:
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
    pic1 = slide3.shapes.add_picture(pic_p1, Inches(0.5), Inches(1.3), width=Inches(new_width1), height=Inches(new_height1))


def create_pres(name):
    global prs
    fullname = 'mabpfolder/' + name
    prs.save(fullname)
    binary_output = BytesIO()
    prs.save(binary_output) 
    print("Done Creating Presentation")
    st.download_button("Download PPT", data = binary_output.getvalue(), file_name = 'test_ppt.pptx')


uploaded_excel = st.file_uploader("Upload Full Excel", type='xlsx')

if uploaded_excel is not None:
    with open('uploaded.xlsx', 'wb') as f:
        f.write(uploaded_excel.getvalue())
        absolute_path = os.path.abspath("uploaded.xlsx") 
        df = load_excel('uploaded.xlsx')
        df_6_1 = clean_sheet(df, 'BS入力')
        create_slide_1(df_6_1)
        df_6_2 = clean_sheet(df, 'PL入力')
        create_slide_2(df_6_2)
        create_slide_3(df_6_2)
        create_pres('test_pres.pptx')
        uploaded_excel = None
        print("Ready to Upload Next File")


